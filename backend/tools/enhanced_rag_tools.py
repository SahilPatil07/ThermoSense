# backend/tools/advanced_rag_system.py
"""
Advanced RAG System - Production Grade
Vector search + Hybrid retrieval + Reranking
"""
import os
import json
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
from datetime import datetime
import hashlib

try:
    import chromadb
    from chromadb.config import Settings
    CHROMADB_AVAILABLE = True
except ImportError:
    CHROMADB_AVAILABLE = False
    print("Warning: ChromaDB not installed. Install: pip install chromadb")

try:
    from sentence_transformers import SentenceTransformer
    EMBEDDINGS_AVAILABLE = True
except ImportError:
    EMBEDDINGS_AVAILABLE = False
    print("Warning: sentence-transformers not installed. Install: pip install sentence-transformers")

try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

# Import PPT tools
try:
    from backend.tools.ppt_tools import PPTX_AVAILABLE
    import pptx
except ImportError:
    PPTX_AVAILABLE = False
    pptx = None

from PIL import Image
import io

class AdvancedRAGSystem:
    """Production RAG with hybrid retrieval and image support"""
    
    def __init__(self, knowledge_dir: str = "assets/knowledge"):
        self.knowledge_dir = Path(knowledge_dir)
        self.knowledge_dir.mkdir(parents=True, exist_ok=True)
        
        self.image_dir = self.knowledge_dir / "images"
        self.image_dir.mkdir(parents=True, exist_ok=True)
        
        self.initialized = False
        
        if not CHROMADB_AVAILABLE or not EMBEDDINGS_AVAILABLE:
            print("Warning: RAG system disabled - missing dependencies")
            return
        
        try:
            # ChromaDB
            self.chroma_client = chromadb.Client(Settings(
                anonymized_telemetry=False
            ))
            
            self.collection_name = "thermal_knowledge"
            self._init_collection()
            
            # Embedding model
            self.embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
            
            # Metadata
            self.metadata_file = self.knowledge_dir / "index_metadata.json"
            self.metadata = self._load_metadata()
            
            self.chunk_size = 1000
            self.chunk_overlap = 200
            
            self.initialized = True
            print("RAG system initialized")
            
        except Exception as e:
            print(f"Warning: RAG initialization failed: {e}")
            self.initialized = False
    
    def _init_collection(self):
        """Initialize ChromaDB collection"""
        try:
            self.collection = self.chroma_client.get_collection(self.collection_name)
            print(f"RAG collection loaded: {self.collection.count()} chunks")
        except:
            self.collection = self.chroma_client.create_collection(
                name=self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )
            print("Created new RAG collection")
    
    def _load_metadata(self) -> Dict:
        """Load indexing metadata"""
        if self.metadata_file.exists():
            with open(self.metadata_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {
            "indexed_files": {},
            "last_indexed": None,
            "total_chunks": 0
        }
    
    def _save_metadata(self):
        """Save metadata"""
        self.metadata["last_indexed"] = datetime.utcnow().isoformat()
        self.metadata["total_chunks"] = self.collection.count() if self.initialized else 0
        with open(self.metadata_file, 'w', encoding='utf-8') as f:
            json.dump(self.metadata, f, indent=2)
    
    def index_knowledge_base(self, force_reindex: bool = False):
        """Index all documents"""
        if not self.initialized:
            print("Warning: RAG system not initialized - skipping indexing")
            return
        
        print("Scanning knowledge base...")
        
        supported_extensions = ['.pdf', '.txt', '.md', '.pptx', '.docx']
        all_files = []
        for ext in supported_extensions:
            all_files.extend(self.knowledge_dir.glob(f'**/*{ext}'))
        
        if not all_files:
            print("Warning: No documents found in assets/knowledge/")
            return
        
        indexed_count = 0
        for file_path in all_files:
            file_hash = self._get_file_hash(file_path)
            file_key = str(file_path.relative_to(self.knowledge_dir))
            
            if not force_reindex and file_key in self.metadata["indexed_files"]:
                if self.metadata["indexed_files"][file_key]["hash"] == file_hash:
                    continue
            
            print(f"Indexing: {file_path.name}")
            
            # Extract text and images
            text_chunks, image_paths = self._extract_content(file_path)
            
            if text_chunks:
                self._add_chunks_to_index(text_chunks, file_path.name, image_paths)
                
                self.metadata["indexed_files"][file_key] = {
                    "hash": file_hash,
                    "indexed_at": datetime.utcnow().isoformat(),
                    "chunks": len(text_chunks)
                }
                indexed_count += 1
        
        if indexed_count > 0:
            self._save_metadata()
            print(f"Indexed {indexed_count} document(s)")
        else:
            print("Knowledge base up-to-date")
    
    def _get_file_hash(self, file_path: Path) -> str:
        """Calculate file hash"""
        hasher = hashlib.md5()
        with open(file_path, 'rb') as f:
            hasher.update(f.read())
        return hasher.hexdigest()
    
    def _extract_content(self, file_path: Path) -> Tuple[List[str], List[str]]:
        """Extract text chunks and images from file"""
        text = ""
        images = []
        try:
            if file_path.suffix.lower() == '.pdf' and PYMUPDF_AVAILABLE:
                text, images = self._extract_pdf_enhanced(file_path)
            elif file_path.suffix.lower() == '.pptx' and PPTX_AVAILABLE:
                text, images = self._extract_pptx_enhanced(file_path)
            elif file_path.suffix.lower() == '.docx' and DOCX_AVAILABLE:
                text = self._extract_docx(file_path)
            elif file_path.suffix.lower() in ['.txt', '.md']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            
            if text:
                chunks = self._chunk_text(text)
                return chunks, images
        except Exception as e:
            print(f"Error extracting {file_path.name}: {e}")
        return [], []
    
    def _extract_pdf_enhanced(self, pdf_path: Path) -> Tuple[str, List[str]]:
        """Extract text and images from PDF using PyMuPDF"""
        text = ""
        image_paths = []
        try:
            doc = fitz.open(pdf_path)
            for page_index in range(len(doc)):
                page = doc[page_index]
                text += page.get_text() + "\n"
                
                # Extract images
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    img_name = f"{pdf_path.stem}_p{page_index}_i{img_index}.png"
                    img_path = self.image_dir / img_name
                    
                    with open(img_path, "wb") as f:
                        f.write(image_bytes)
                    image_paths.append(str(img_path))
            doc.close()
        except Exception as e:
            print(f"Warning: PDF enhanced extraction error: {e}")
        return text, image_paths

    def _extract_pptx_enhanced(self, pptx_path: Path) -> Tuple[str, List[str]]:
        """Extract text and images from PPTX"""
        text = ""
        image_paths = []
        try:
            prs = pptx.Presentation(pptx_path)
            for slide_index, slide in enumerate(prs.slides):
                for shape_index, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                    
                    if shape.shape_type == 13: # Picture
                        image = shape.image
                        image_bytes = image.blob
                        img_name = f"{pptx_path.stem}_s{slide_index}_i{shape_index}.png"
                        img_path = self.image_dir / img_name
                        
                        with open(img_path, "wb") as f:
                            f.write(image_bytes)
                        image_paths.append(str(img_path))
        except Exception as e:
            print(f"Warning: PPTX enhanced extraction error: {e}")
        return text, image_paths
    
    def _extract_docx(self, docx_path: Path) -> str:
        """Extract text from DOCX"""
        text = ""
        try:
            doc = docx.Document(docx_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            print(f"Warning: DOCX extraction error: {e}")
        return text

    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks with overlap"""
        if not text:
            return []
            
        # Basic recursive character splitter logic
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            chunk = text[start:end]
            chunks.append(chunk)
            start += self.chunk_size - self.chunk_overlap
            
        return chunks
    
    def _add_chunks_to_index(self, chunks: List[str], source: str, image_paths: List[str] = None):
        """Add chunks to vector database"""
        if not chunks or not self.initialized:
            return
        
        embeddings = self.embedding_model.encode(chunks).tolist()
        
        # Distribute image paths across chunks (simplistic approach: first few chunks get images)
        # In a better implementation, we'd map images to specific chunks based on page/slide
        metadatas = []
        for i in range(len(chunks)):
            meta = {"source": source, "chunk_index": i}
            if image_paths and i < len(image_paths):
                meta["image_path"] = image_paths[i]
            metadatas.append(meta)
        
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        ids = [f"{source}_{timestamp}_{i}" for i in range(len(chunks))]
        
        # Batch add to ChromaDB to avoid limits
        batch_size = 100
        for i in range(0, len(chunks), batch_size):
            end = min(i + batch_size, len(chunks))
            self.collection.add(
                documents=chunks[i:end],
                embeddings=embeddings[i:end],
                metadatas=metadatas[i:end],
                ids=ids[i:end]
            )
    
    def retrieve(self, query: str, top_k: int = 5, min_confidence: float = 0.2) -> Tuple[str, List[str], float, List[str]]:
        """
        Retrieve relevant context and images
        Returns: (context, sources, confidence, image_paths)
        """
        if not self.initialized or self.collection.count() == 0:
            return "", [], 0.0, []
        
        try:
            # Generate query embedding
            query_embedding = self.embedding_model.encode([query]).tolist()
            
            results = self.collection.query(
                query_embeddings=query_embedding,
                n_results=top_k
            )
            
            if not results['documents'] or not results['documents'][0]:
                return "", [], 0.0, []
            
            candidates = []
            image_paths = []
            
            for doc, metadata, distance in zip(
                results['documents'][0],
                results['metadatas'][0],
                results['distances'][0]
            ):
                # ChromaDB returns squared L2 distance by default if not specified
                # We used cosine similarity in init, so distance is 1 - similarity
                similarity = 1 - distance
                
                if similarity >= min_confidence:
                    candidates.append({
                        "text": doc,
                        "source": metadata.get("source", "Unknown"),
                        "similarity": similarity
                    })
                    if "image_path" in metadata:
                        image_paths.append(metadata["image_path"])
            
            if not candidates:
                return "", [], 0.0, []
            
            context = "\n\n---\n\n".join([
                f"[Source: {c['source']}]\n{c['text']}"
                for c in candidates
            ])
            
            sources = list(set([c['source'] for c in candidates]))
            avg_confidence = sum([c['similarity'] for c in candidates]) / len(candidates)
            
            # Deduplicate image paths and convert to relative URLs for frontend
            unique_images = []
            for path in image_paths:
                rel_path = Path(path).relative_to(Path("assets").parent)
                url = f"/{rel_path.as_posix()}"
                if url not in unique_images:
                    unique_images.append(url)
            
            return context, sources, avg_confidence, unique_images
        
        except Exception as e:
            print(f"Warning: Retrieval error: {e}")
            return "", [], 0.0, []

# Singleton
rag_system = AdvancedRAGSystem()

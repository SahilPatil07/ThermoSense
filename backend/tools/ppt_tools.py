"""
PowerPoint Slide Extraction Tools
"""
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None
    Inches = None

from pathlib import Path
from typing import List, Optional, Dict, Any
import io
import logging
import uuid

logger = logging.getLogger(__name__)


class PPTExtractor:
    """Extract slides from PowerPoint presentations"""
    
    def __init__(self):
        pass
    
    def get_slide_info(self, ppt_path: str) -> List[Dict[str, Any]]:
        """Get titles and numbers for all slides"""
        if not PPTX_AVAILABLE:
            return []
            
        try:
            if ppt_path.lower().endswith('.ppt'):
                # For .ppt, we must use COM as python-pptx only supports .pptx
                return self._get_slide_info_com(ppt_path)
            
            prs = Presentation(ppt_path)
            info = []
            for i, slide in enumerate(prs.slides):
                info.append({
                    "number": i + 1,
                    "title": self._get_slide_title(slide)
                })
            return info
        except Exception as e:
            logger.error(f"Error getting slide info: {e}")
            return []

    def _get_slide_info_com(self, ppt_path: str) -> List[Dict[str, Any]]:
        """Get slide info using COM (for .ppt files)"""
        info = []
        ppt_app = None
        presentation = None
        try:
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            abs_path = str(Path(ppt_path).absolute())
            presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, Untitled=False, WithWindow=False)
            
            for i in range(1, presentation.Slides.Count + 1):
                slide = presentation.Slides(i)
                title = "Untitled Slide"
                try:
                    if slide.Shapes.HasTitle:
                        title = slide.Shapes.Title.TextFrame.TextRange.Text
                except:
                    pass
                info.append({"number": i, "title": title})
            return info
        except Exception as e:
            logger.error(f"COM slide info failed: {e}")
            return []
        finally:
            if presentation: presentation.Close()
            if ppt_app: ppt_app.Quit()
            try: pythoncom.CoUninitialize()
            except: pass

    def export_slides_as_images(self, ppt_path: str, slide_numbers: List[int], output_dir: Path) -> List[str]:
        """
        Export specific slides as high-quality images using PowerPoint COM.
        Falls back to text extraction if COM fails.
        """
        image_paths = []
        ppt_app = None
        presentation = None
        
        try:
            try:
                import win32com.client
                import pythoncom
            except ImportError:
                logger.warning("pywin32 (win32com) not installed. PPT image export requires pywin32. Using text-only fallback.")
                logger.info("To enable full PPT image export on Windows: pip install pywin32 && python venv/Scripts/pywin32_postinstall.py -install")
                return []  # Fallback will be handled in extract_slides
            
            # Initialize COM
            pythoncom.CoInitialize()
            
            try:
                ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            except Exception as e:
                logger.error(f"Failed to start PowerPoint application via COM: {e}")
                logger.info("Ensure Microsoft PowerPoint is installed on this Windows machine.")
                return []

            abs_ppt_path = str(Path(ppt_path).absolute())
            try:
                presentation = ppt_app.Presentations.Open(abs_ppt_path, ReadOnly=True, Untitled=False, WithWindow=False)
            except Exception as e:
                logger.error(f"Failed to open presentation via COM: {e}")
                return []
            
            output_dir.mkdir(parents=True, exist_ok=True)
            
            for slide_num in slide_numbers:
                if 1 <= slide_num <= presentation.Slides.Count:
                    filename = f"slide_{slide_num}_{uuid.uuid4().hex[:8]}.png"
                    output_path = output_dir / filename
                    
                    # Export slide as PNG
                    # 2 = ppShapeFormatPNG
                    try:
                        presentation.Slides(slide_num).Export(str(output_path), "PNG")
                        image_paths.append(str(output_path))
                        logger.info(f"✓ Exported slide {slide_num} to {filename}")
                    except Exception as e:
                        logger.error(f"Failed to export slide {slide_num}: {e}")
                else:
                    logger.warning(f"Slide {slide_num} out of range (Total: {presentation.Slides.Count})")
            
            return image_paths
            
        except Exception as e:
            logger.error(f"Unexpected COM Export error: {e}")
            return []
        finally:
            if presentation:
                try:
                    presentation.Close()
                except:
                    pass
            if ppt_app:
                try:
                    ppt_app.Quit()
                except:
                    pass
            try:
                pythoncom.CoUninitialize()
            except:
                pass

    def extract_slides_text_fallback(self, ppt_path: str, slide_numbers: List[int], output_dir: Optional[Path] = None) -> List[Dict[str, Any]]:
        """
        Fallback: Extract slides as text/metadata when image export fails.
        Uses python-pptx to extract all text content, notes, and structure.
        """
        if not PPTX_AVAILABLE:
            logger.error("python-pptx not available for text extraction")
            return []
        
        try:
            prs = Presentation(ppt_path)
            extracted_slides = []
            
            logger.info(f"Using text-only extraction fallback for {len(slide_numbers)} slides")
            
            for slide_num in slide_numbers:
                idx = slide_num - 1
                
                if 0 <= idx < len(prs.slides):
                    slide = prs.slides[idx]
                    
                    # Extract all text and create a structured representation
                    text_content = self._get_slide_text(slide)
                    title = self._get_slide_title(slide)
                    
                    # Create a simple text file as fallback "image"
                    text_file_path = None
                    if output_dir:
                        output_dir.mkdir(parents=True, exist_ok=True)
                        text_filename = f"slide_{slide_num}_text_{uuid.uuid4().hex[:8]}.txt"
                        text_file_path = output_dir / text_filename
                        
                        with open(text_file_path, 'w', encoding='utf-8') as f:
                            f.write(f"Slide {slide_num}: {title}\n")
                            f.write("=" * 50 + "\n\n")
                            f.write(text_content)
                        
                        logger.info(f"✓ Extracted text from slide {slide_num} to {text_filename}")
                    
                    slide_data = {
                        'slide_number': slide_num,
                        'title': title,
                        'text_content': text_content,
                        'image_path': str(text_file_path) if text_file_path else None,
                        'shapes_count': len(slide.shapes),
                        'extraction_mode': 'text_only'
                    }
                    extracted_slides.append(slide_data)
                else:
                    logger.warning(f"Slide {slide_num} not found in {ppt_path}")
            
            return extracted_slides
            
        except Exception as e:
            logger.error(f"Error in text fallback extraction from {ppt_path}: {e}")
            return []
    
    def extract_slides(self, ppt_path: str, slide_numbers: List[int], output_dir: Optional[Path] = None) -> List[Dict[str, Any]]:
        """
        Extract specific slides from PowerPoint.
        Tries image export first, falls back to text-only extraction.
        """
        if not PPTX_AVAILABLE:
            logger.warning("python-pptx not installed, cannot extract slides")
            return []
            
        try:
            prs = Presentation(ppt_path)
            extracted_slides = []
            
            # Try to export images if output_dir is provided
            image_map = {}
            com_available = False
            
            if output_dir:
                exported_images = self.export_slides_as_images(ppt_path, slide_numbers, output_dir)
                if exported_images:
                    # COM export succeeded
                    com_available = True
                    for i, slide_num in enumerate(slide_numbers):
                        if i < len(exported_images):
                            image_map[slide_num] = exported_images[i]
                    logger.info(f"Successfully exported {len(exported_images)} slides as images")
                else:
                    # COM export failed, use text fallback
                    logger.info("Image export unavailable, using text-only extraction")
                    return self.extract_slides_text_fallback(ppt_path, slide_numbers, output_dir)

            # Process slides with image paths
            for slide_num in slide_numbers:
                idx = slide_num - 1
                
                if 0 <= idx < len(prs.slides):
                    slide = prs.slides[idx]
                    
                    slide_data = {
                        'slide_number': slide_num,
                        'title': self._get_slide_title(slide),
                        'text_content': self._get_slide_text(slide),
                        'image_path': image_map.get(slide_num),
                        'shapes_count': len(slide.shapes),
                        'extraction_mode': 'full' if com_available else 'metadata_only'
                    }
                    extracted_slides.append(slide_data)
                else:
                    logger.warning(f"Slide {slide_num} not found in {ppt_path}")
            
            return extracted_slides
            
        except Exception as e:
            logger.error(f"Error extracting slides from {ppt_path}: {e}")
            return []
    
    def _get_slide_title(self, slide) -> str:
        """Extract title from slide"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text
            # Fallback: Check for first textbox
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    return shape.text.strip().split('\n')[0][:50]
            return "Untitled Slide"
        except:
            return "Untitled Slide"
    
    def _get_slide_text(self, slide) -> str:
        """Extract all text from slide"""
        try:
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            return ""
    
    def extract_images_from_slide(self, slide, output_dir: Path, prefix: str) -> List[str]:
        """Extract embedded images from slide (raw blobs)"""
        image_paths = []
        try:
            for i, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # Picture
                    image = shape.image
                    image_bytes = image.blob
                    ext = image.ext
                    
                    filename = f"{prefix}_img_{i}.{ext}"
                    path = output_dir / filename
                    
                    with open(path, 'wb') as f:
                        f.write(image_bytes)
                    
                    image_paths.append(str(path))
        except Exception as e:
            logger.error(f"Image extraction error: {e}")
        return image_paths

    def _get_slide_images(self, slide) -> List[str]:
        """Get image descriptions (legacy)"""
        return []
    
    def copy_slides_to_document(self, ppt_path: str, slide_numbers: List[int], 
                               output_pptx: str) -> bool:
        """
        Create new PPTX with only selected slides
        """
        if not PPTX_AVAILABLE:
            return False
            
        try:
            source_prs = Presentation(ppt_path)
            new_prs = Presentation()
            
            # Remove default blank slide
            if len(new_prs.slides) > 0:
                for i in range(len(new_prs.slides) - 1, -1, -1):
                    rId = new_prs.slides._sldIdLst[i].rId
                    new_prs.part.drop_rel(rId)
                    del new_prs.slides._sldIdLst[i]
            
            # Copy selected slides
            for slide_num in slide_numbers:
                idx = slide_num - 1
                
                if 0 <= idx < len(source_prs.slides):
                    source_slide = source_prs.slides[idx]
                    
                    # Copy slide layout
                    slide_layout = new_prs.slide_layouts[0]
                    new_slide = new_prs.slides.add_slide(slide_layout)
                    
                    # Copy shapes from source
                    for shape in source_slide.shapes:
                        self._copy_shape(shape, new_slide)
            
            # Save new presentation
            new_prs.save(output_pptx)
            return True
            
        except Exception as e:
            logger.error(f"Error copying slides: {e}")
            return False
    
    def _copy_shape(self, source_shape, target_slide):
        """Copy a shape to target slide"""
        try:
            if hasattr(source_shape, "text"):
                el = target_slide.shapes.add_textbox(
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height
                )
                el.text = source_shape.text
        except Exception as e:
            logger.error(f"Error copying shape: {e}")
    
    def get_slide_count(self, ppt_path: str) -> int:
        """Get total number of slides in presentation"""
        if not PPTX_AVAILABLE:
            return 0
            
        try:
            if ppt_path.lower().endswith('.ppt'):
                # Use COM for .ppt
                return len(self._get_slide_info_com(ppt_path))
            prs = Presentation(ppt_path)
            return len(prs.slides)
        except Exception as e:
            logger.error(f"Error getting slide count: {e}")
            return 0

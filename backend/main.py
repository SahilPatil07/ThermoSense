import re
import os
import sys
import json
import uuid
import logging
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional
from contextlib import asynccontextmanager
from collections import deque
import asyncio
import functools

# Add parent directory to path for imports when running directly
if __name__ == "__main__":
    parent_dir = Path(__file__).resolve().parent.parent
    if str(parent_dir) not in sys.path:
        sys.path.insert(0, str(parent_dir))

import pandas as pd
import numpy as np
import uvicorn
from fastapi import FastAPI, UploadFile, File, Form, Request, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field


from backend.tools.storage import SessionStorage
from backend.tools.chart_tools import ChartGenerator
from backend.tools.chart_tools_plotly import PlotlyChartGenerator
from backend.tools.ppt_tools import PPTExtractor
from backend.tools.llm_orchestrator import LLMOrchestrator
from backend.tools.excel_tools import read_excel_table, list_sheetnames
from backend.tools.advanced_analytics import ThermalAnalytics
from backend.tools.enhanced_llm_orchestrator import get_enhanced_orchestrator
from backend.tools.data_context import get_context_manager
from backend.tools.sensor_harvester import SensorHarvester
from backend.tools.comparative_analyzer import ComparativeAnalyzer
from backend.tools.utils import detect_time_column, get_numeric_columns

# Observability imports
from backend.observability.logger import setup_logger, set_context
from backend.observability.metrics import metrics, time_execution
from backend.observability.tracing import trace_span, tracer

# API imports
from backend.api import agent_endpoints, health_endpoints, signal_endpoints





# Logging configuration
logger = setup_logger("backend")
logger.info("Structured logging initialized")


# ========== STEP 1: CREATE DIRECTORIES FIRST ==========
def ensure_directories():
    """Create all required directories"""
    dirs = [
        "static",
        "workspace",
        "workspace/approved_charts",
        "knowledge_base",
        "assets",
        "assets/knowledge",
        "backend/tools"
    ]
    for directory in dirs:
        Path(directory).mkdir(parents=True, exist_ok=True)

# Call this BEFORE any imports
ensure_directories()
logger.info("Directories ready")



# Initialize core components
storage = SessionStorage(workspace_dir="workspace")
chart_generator = ChartGenerator(dpi=300, figsize=(14, 8))
plotly_generator = PlotlyChartGenerator()
ppt_extractor = PPTExtractor()

# Initialize context management
context_manager = get_context_manager()
sensor_harvester = SensorHarvester()
comparative_analyzer = ComparativeAnalyzer()

# LLM Client (Ollama)
llm_client = None
try:
    from openai import OpenAI
    
    logger.info(f"OpenAI Class: {OpenAI}")
    logger.info(f"OpenAI Module: {sys.modules.get('openai')}")
    
    # Check for proxy env vars
    proxy_vars = {k: v for k, v in os.environ.items() if 'PROXY' in k.upper()}
    logger.info(f"Proxy Env Vars: {proxy_vars}")

    # Attempt initialization
    logger.info("Initializing OpenAI client...")
    llm_client = OpenAI(base_url="http://localhost:11434/v1", api_key="ollama")
    
    orchestrator = LLMOrchestrator(llm_client=llm_client)
    enhanced_orchestrator = get_enhanced_orchestrator(llm_client)
    analytics_engine = ThermalAnalytics(llm_client=llm_client)
    logger.info("Ollama LLM connected successfully")
except Exception as e:
    logger.error(f"LLM initialization failed: {e}")
    import traceback
    logger.error(traceback.format_exc())
    
    # Fallback
    orchestrator = LLMOrchestrator(llm_client=None)
    enhanced_orchestrator = get_enhanced_orchestrator(None)
    analytics_engine = ThermalAnalytics(llm_client=None)


# ========== STEP 3: LIFESPAN MANAGEMENT ==========
@asynccontextmanager
async def lifespan(app: FastAPI):
    global agent_registry, agent_planner, agent_router, agent_verifier
    
    logger.info("ThermoSense AI Starting...")
    
    # Try to load RAG (optional)
    try:
        from backend.tools.enhanced_rag_tools import rag_system
        logger.info("Indexing knowledge base...")
        rag_system.index_knowledge_base(force_reindex=False)
        logger.info("RAG system ready")
    except ImportError:
        logger.info("RAG system not installed (optional feature)")
    except Exception as e:
        logger.warning(f"RAG initialization failed: {e}")
    
    # Initialize agent platform
    try:
        from backend.agent.setup import initialize_agent_platform
        from backend.tools.report_generator import ThermalReportGenerator
        
        report_gen = ThermalReportGenerator(llm_client=llm_client)
        
        agent_registry, agent_planner, agent_router, agent_verifier, _ = initialize_agent_platform(
            storage=storage,
            sensor_harvester=sensor_harvester,
            plotly_generator=plotly_generator,
            comparative_analyzer=comparative_analyzer,
            report_generator=report_gen,
            llm_client=llm_client
        )
        
        # Set components in agent endpoints
        from backend.api import agent_endpoints
        agent_endpoints.set_agent_components(agent_planner, agent_router, agent_registry)

        # Set planner in enhanced orchestrator for unified chat
        enhanced_orchestrator.planner = agent_planner
        
        logger.info("Agent platform initialized and integrated with orchestrator")
    except Exception as e:
        logger.warning(f"Agent platform initialization failed: {e}")
    
    logger.info("System Ready")
    yield
    logger.info("Shutting down")


# ========== FASTAPI APP ==========
app = FastAPI(title="ThermoSense AI", lifespan=lifespan)

# Include routers
try:
    from backend.api.task_endpoints import router as task_router
    from backend.api import signal_endpoints
    app.include_router(task_router)
    app.include_router(agent_endpoints.router)
    app.include_router(signal_endpoints.router)
    app.include_router(health_endpoints.router)
    app.include_router(health_endpoints.health_router)
    logger.info("API routers mounted successfully")
except Exception as e:
    logger.error(f"Router mounting failed: {e}")

app.add_middleware(


    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"]
)

# Mount static files
app.mount("/static", StaticFiles(directory="static"), name="static")
app.mount("/workspace", StaticFiles(directory="workspace"), name="workspace")

# (Routers moved up)





# ========== REQUEST MODELS ==========
class ChatRequest(BaseModel):
    message: str
    session_id: str = "default"

class ColumnSelectionRequest(BaseModel):
    session_id: str
    filename: str  # Changed from 'file' to match other models
    x_column: Optional[str]
    y_columns: List[str]

class FeedbackRequest(BaseModel):
    session_id: str
    chart_id: str
    feedback: str

class ChartRequest(BaseModel):
    session_id: str
    filename: str  # Changed from 'file' to match frontend
    x_column: str
    y_columns: List[str]
    chart_type: str = "line"
    user_query: Optional[str] = None
    sheet_name: Optional[str] = None

class HeatmapRequest(BaseModel):
    session_id: str
    filename: str  # Changed from 'file' to match frontend
    columns: Optional[List[str]] = None

class ExtractionRequest(BaseModel):
    session_id: str
    filenames: List[str]
    sensors: List[str] = []
    strict: bool = False
    structured: bool = True
    sheet_name: Optional[str] = None
    columns: Optional[List[str]] = None
    section: Optional[str] = None

class InferenceRequest(BaseModel):
    session_id: str
    content: str
    chart_id: Optional[str] = None
    section: Optional[str] = None


# ========== HELPER FUNCTIONS ==========

def fallback_response(user_message: str, uploads: List[str], approved_count: int) -> str:
    """Intelligent fallback when LLM unavailable"""
    msg_lower = user_message.lower()
    
    # Greetings
    if any(word in msg_lower for word in ['hi', 'hey', 'hello', 'greetings']):
        return f"""Hello! I'm ThermoSense AI, your thermal analysis assistant.

**I can help you:**
• Analyze thermal data ({len(uploads)} file{"s" if len(uploads) != 1 else ""} uploaded)
• Create visualizations (charts, heatmaps)
• Answer technical questions
• Generate reports ({approved_count} charts approved)

**Try asking:**
- "What is heat transfer?"
- "Explain thermal conductivity"
- "Create a chart" (select columns)
- "Generate report"

What would you like to know?"""
    
    # Technical questions
    if any(word in msg_lower for word in ['what is', 'define', 'explain']):
        if 'mass' in msg_lower and 'flow' in msg_lower:
            return """**Mass Flow Rate** is the amount of mass flowing through a cross-section per unit time.

**Formula:** ṁ = ρ × V × A

Where:
• ṁ = mass flow rate (kg/s)
• ρ = fluid density (kg/m³)
• V = fluid velocity (m/s)
• A = cross-sectional area (m²)

**Automotive applications:**
- Coolant flow through radiator
- Air flow through heat exchanger
- Fuel flow in engine

**Example:** If coolant (ρ=1000 kg/m³) flows at 2 m/s through 0.01 m² pipe:
ṁ = 1000 × 2 × 0.01 = **20 kg/s**

Would you like to analyze your thermal data?"""
        
        return "I can explain thermal engineering concepts! Ask about heat transfer, thermal conductivity, mass flow rate, or cooling systems."
    
    # Data analysis
    if uploads:
        return f"""You have **{len(uploads)} file(s)**: {', '.join(uploads)}

**To analyze:**
1. Click a file in sidebar
2. Select columns (X and Y axes)
3. I'll generate a chart
4. Click 👍 to approve for reports

**Or ask:**
- "Show temperature trends"
- "Create a heatmap"
- "Generate report"

What would you like?"""
    
    return """I'm ThermoSense AI - your thermal expert!

**Upload files** (.csv, .xlsx) to get started
**Ask technical questions** about thermal systems
**Generate reports** with visualizations

What can I help you with?"""


# ========== API ENDPOINTS ==========

@app.get("/")
async def root():
    """Serve main UI"""
    return FileResponse("static/index.html")

@app.post("/api/upload")
async def upload_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...), 
    session_id: str = Form("default")
):
    try:
        session_dir = storage.get_session_dir(session_id)
        
        # Ensure directory exists
        session_dir.mkdir(parents=True, exist_ok=True)
        
        file_path = session_dir / file.filename
        
        logger.info(f"Saving file to: {file_path}")
        
        # Validate file extension
        allowed_extensions = {'.csv', '.xlsx', '.xls', '.pptx'}
        file_ext = Path(file.filename).suffix.lower()
        
        if file_ext not in allowed_extensions:
            return JSONResponse({
                "success": False,
                "error": f"Invalid file type: {file_ext}. Allowed: {', '.join(allowed_extensions)}"
            }, status_code=400)

        # Save file
        with open(str(file_path), "wb") as f:  # Convert Path to string
            content = await file.read()
            f.write(content)
        
        # Verify file was saved
        if not file_path.exists():
            raise Exception(f"File save failed: {file_path}")
        
        storage.add_upload(session_id, file.filename)
        logger.info(f"Uploaded: {file.filename} ({file_path.stat().st_size} bytes)")
        
        columns = []
        time_column = None
        numeric_columns = []
        
        try:
            if file.filename.lower().endswith(('.xlsx', '.xls', '.csv')):
                columns, time_column, numeric_columns = sensor_harvester.get_file_columns(str(file_path))
                
                logger.info(f"Extracted {len(columns)} valid columns, time: {time_column}")
                
                # Store file context for intelligent querying
                if file.filename.endswith('.csv'):
                    df = pd.read_csv(str(file_path))
                else:
                    df = pd.read_excel(str(file_path))
                
                column_types = {col: str(df[col].dtype) for col in columns if col in df.columns}
                data_preview = {
                    'row_count': len(df),
                    'sample_values': {col: df[col].head(3).tolist() for col in columns[:5] if col in df.columns}
                }
                context_manager.store_file_context(
                    session_id=session_id,
                    filename=file.filename,
                    columns=columns,
                    row_count=len(df),
                    column_types=column_types,
                    data_preview=data_preview
                )
                
                # Trigger background analysis for data files (Disabled as per user request)
                # background_tasks.add_task(run_proactive_analysis, session_id, file.filename, file_path)

            elif file.filename.lower().endswith(('.pptx', '.ppt')):
                slide_info = ppt_extractor.get_slide_info(str(file_path))
                slide_count = len(slide_info)
                logger.info(f"PPTX uploaded: {slide_count} slides")
                return JSONResponse({
                    "success": True,
                    "filename": file.filename,
                    "is_ppt": True,
                    "slide_count": slide_count,
                    "slide_info": slide_info
                })
        except Exception as e:
            logger.error(f"Metadata extraction error: {e}")
        
        sheets = []
        if file.filename.lower().endswith(('.xlsx', '.xls')):
             try:
                 sheets = list_sheetnames(str(file_path))
             except:
                 pass

        return JSONResponse({
            "success": True,
            "filename": file.filename,
            "columns": columns,
            "time_column": time_column,
            "numeric_columns": numeric_columns,
            "sheets": sheets
        })
        
    except Exception as e:
        logger.exception("Upload error")
        return JSONResponse({
            "success": False, 
            "error": str(e)
        }, status_code=500)

async def run_proactive_analysis(session_id: str, filename: str, file_path: Path):
    """
    Background task to perform proactive analysis and generate insights.
    """
    try:
        logger.info(f"🚀 Starting proactive analysis for {filename} in session {session_id}")
        
        # Load data
        if filename.endswith('.csv'):
            df = pd.read_csv(str(file_path))
        else:
            sheets = list_sheetnames(str(file_path))
            df = read_excel_table(str(file_path), sheets[0])
            
        # Detect columns
        cols, time_col, num_cols = sensor_harvester.get_file_columns(str(file_path))
        
        if not num_cols:
            logger.warning(f"No numeric columns found in {filename} for proactive analysis")
            return
            
        # Run analysis
        # Optimization: Limit to top 3 numeric columns to save time
        analysis = analytics_engine.comprehensive_analysis(
            df=df,
            x_col=time_col or df.columns[0],
            y_cols=num_cols[:3] # Reduced from 5 to 3 for speed
        )
        
        # Store analysis
        storage.set_analysis(session_id, filename, analysis)
        
        # Generate insight
        insight_text = enhanced_orchestrator.get_proactive_insight(analysis, filename)
        
        # Store insight
        storage.add_proactive_insight(session_id, {
            "filename": filename,
            "content": insight_text,
            "type": "proactive_insight"
        })
        
        logger.info(f"✅ Proactive analysis complete for {filename}")
        
    except Exception as e:
        logger.error(f"Proactive analysis failed for {filename}: {e}")

@app.get("/api/analysis/proactive")
async def get_proactive_insight(session_id: str = "default"):
    """
    Fetch the latest unread proactive insight for the session.
    """
    insight = storage.get_latest_proactive_insight(session_id)
    if insight:
        return JSONResponse({"success": True, "insight": insight})
    return JSONResponse({"success": False, "message": "No new insights"})

class RecommendRequest(BaseModel):
    session_id: str
    y_columns: List[str]

@app.post("/api/chart/recommend")
async def recommend_chart(req: RecommendRequest):
    """
    Consolidated chart recommendation endpoint
    - Uses historical patterns (global)
    - Uses session history (local)
    - Uses rule-based logic
    - Uses LLM (optional)
    """
    try:
        recommendations = []
        seen_types = set()
        
        # 1. Get historical patterns (global)
        hist_patterns = storage.get_recommendations_for_signals(req.y_columns)
        for r in hist_patterns:
            if r not in seen_types:
                recommendations.append({
                    "type": r, 
                    "reason": "Based on global successful patterns", 
                    "confidence": 0.9, 
                    "source": "history_global"
                })
                seen_types.add(r)
        
        # 2. Check session history (local)
        memory = storage.load_memory(req.session_id)
        session_history = memory.get("chart_history", [])
        session_prefs = {}
        for chart in session_history:
            if chart.get("approved"):
                ctype = chart.get("chart_type", "line")
                session_prefs[ctype] = session_prefs.get(ctype, 0) + 1
        
        if session_prefs:
            top_type = max(session_prefs, key=session_prefs.get)
            if top_type not in seen_types:
                recommendations.append({
                    "type": top_type,
                    "reason": "You previously approved this type in this session",
                    "confidence": 0.8,
                    "source": "history_session"
                })
                seen_types.add(top_type)
        
        # 3. Rule-based recommendations
        y_cols_lower = [col.lower() for col in req.y_columns]
        
        if len(req.y_columns) > 5 and "heatmap" not in seen_types:
            recommendations.append({
                "type": "heatmap",
                "reason": "Heatmaps are great for many sensors",
                "confidence": 0.7,
                "source": "rules"
            })
            seen_types.add("heatmap")
            
        if any("temp" in col for col in y_cols_lower) and "line" not in seen_types:
            recommendations.append({
                "type": "line",
                "reason": "Ideal for temperature trends over time",
                "confidence": 0.7,
                "source": "rules"
            })
            seen_types.add("line")
            
        if any(kw in col for col in y_cols_lower for kw in ["power", "current"]) and "bar" not in seen_types:
            recommendations.append({
                "type": "bar",
                "reason": "Good for power comparison across states",
                "confidence": 0.6,
                "source": "rules"
            })
            seen_types.add("bar")
            
        # 4. LLM recommendations (if available)
        try:
            orchestrator = get_enhanced_orchestrator()
            llm_recs = orchestrator.get_chart_recommendations(req.y_columns)
            for r in llm_recs:
                if r["type"] not in seen_types:
                    recommendations.append({**r, "source": "llm"})
                    seen_types.add(r["type"])
        except Exception as e:
            logger.warning("LLM recommendation failed: " + str(e))
            
        # 5. Default if nothing else
        if not recommendations:
            recommendations.append({
                "type": "line",
                "reason": "Standard for sensor data",
                "confidence": 0.5,
                "source": "default"
            })

        return JSONResponse({
            "success": True, 
            "recommendations": recommendations[:3]
        })
    except Exception as e:
        logger.error("Recommendation error: " + str(e))
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.get("/api/uploads")
async def get_uploads(session_id: str = "default"):
    """Get all uploaded files for session"""
    # Get session-specific files
    uploads = storage.get_uploads(session_id)
    files_info = []
    
    for fname in uploads:
        # Find the file path (it could be in any session dir)
        fpath = storage.find_file(fname)
        if fpath and fpath.exists():
            columns = []
            time_column = None
            numeric_columns = []
            
            try:
                if fname.lower().endswith(('.xlsx', '.xls', '.csv')):
                    columns, time_column, numeric_columns = sensor_harvester.get_file_columns(str(fpath))
            except:
                pass
            
            files_info.append({
                "filename": fname,
                "size": fpath.stat().st_size,
                "columns": columns,
                "time_column": time_column,
                "numeric_columns": numeric_columns
            })
    
    return JSONResponse({"files": files_info})

@app.get("/api/file/sheets")
async def get_file_sheets(filename: str, session_id: str = "default"):
    """Get sheet names for an Excel file"""
    try:
        logger.info(f"Listing sheets for: {filename} in session: {session_id}")
        # Prioritize session directory
        session_dir = storage.get_session_dir(session_id)
        fpath = session_dir / filename
        
        if not fpath.exists():
            logger.info(f"File not in session dir, searching workspace: {filename}")
            # Fallback to global search
            fpath = storage.find_file(filename)
            
        if not fpath or not fpath.exists():
            logger.warning(f"File not found: {filename}")
            return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
        
        logger.info(f"Found file at: {fpath}")
        if filename.lower().endswith(('.xlsx', '.xls')):
            sheets = list_sheetnames(str(fpath))
            return JSONResponse({"success": True, "sheets": sheets})
        return JSONResponse({"success": False, "error": "Not an Excel file"}, status_code=400)
    except Exception as e:
        logger.error(f"Error listing sheets: {e}")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.post("/api/file/sheets")
async def post_file_sheets(request: Request):
    """Get sheet names for an Excel file (POST version for robust filename handling)"""
    try:
        data = await request.json()
        filename = data.get("filename")
        session_id = data.get("session_id", "default")
        
        if not filename:
            return JSONResponse({"success": False, "error": "Filename required"}, status_code=400)
            
        return await get_file_sheets(filename, session_id)
    except Exception as e:
        logger.error(f"Error in post_file_sheets: {e}")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.get("/api/file/columns")
async def get_file_columns(filename: str, sheet_name: str, session_id: str = "default"):
    """Get column names for a specific sheet in an Excel file"""
    try:
        logger.info(f"Listing columns for: {filename} [{sheet_name}] in session: {session_id}")
        # Prioritize session directory
        session_dir = storage.get_session_dir(session_id)
        fpath = session_dir / filename
        
        if not fpath.exists():
            logger.info(f"File not in session dir, searching workspace: {filename}")
            # Fallback to global search
            fpath = storage.find_file(filename)
            
        if not fpath or not fpath.exists():
            logger.warning(f"File not found: {filename}")
            return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
        
        logger.info(f"Found file at: {fpath}")
        try:
            df = read_excel_table(str(fpath), sheet_name)
            columns = list(df.columns)
            return JSONResponse({"success": True, "columns": columns})
        except ValueError as e:
            # Worksheet not found error
            logger.error(f"Worksheet error: {e}")
            error_str = str(e)
            if "not found" in error_str.lower():
                return JSONResponse({
                    "success": False,
                    "error": f"Worksheet '{sheet_name}' not found or invalid.",
                    "suggestion": "Please check the sheet name and try again.",
                    "technical_details": error_str
                }, status_code=404)
            else:
                # Other ValueErrors (e.g. data issues)
                return JSONResponse({
                    "success": False,
                    "error": "Data error in Excel sheet.",
                    "suggestion": "Check if the sheet has merged cells or special formatting.",
                    "technical_details": error_str
                }, status_code=500)
        except Exception as e:
            # Other errors (openpyxl, pandas, etc.)
            logger.error(f"Error listing columns: {e}")
            error_msg = str(e)
            
            # Provide helpful suggestions based on error type
            if "calculate_dimension" in error_msg or "wildcard" in error_msg:
                suggestion = "This sheet may have invalid cell references or special formatting. Try opening and re-saving the file."
            elif "openpyxl" in error_msg:
                suggestion = "The Excel file may be corrupted. Try opening it in Excel and saving a clean copy."
            else:
                suggestion = "Please check if the file is a valid Excel file and not password-protected."
            
            return JSONResponse({
                "success": False,
                "error": "Cannot process this Excel sheet.",
                "suggestion": suggestion,
                "technical_details": error_msg
            }, status_code=500)
    except Exception as outer_e:
        logger.exception("Unexpected error in get_file_columns")
        return JSONResponse({"success": False, "error": "Internal server error"}, status_code=500)

@app.post("/api/file/columns")
async def post_file_columns(request: Request):
    """Get column names for a specific sheet in an Excel file (POST version)"""
    try:
        data = await request.json()
        filename = data.get("filename")
        sheet_name = data.get("sheet_name")
        session_id = data.get("session_id", "default")
        
        if not filename or not sheet_name:
            return JSONResponse({"success": False, "error": "Filename and sheet_name required"}, status_code=400)
            
        return await get_file_columns(filename, sheet_name, session_id)
    except Exception as e:
        logger.error(f"Error in post_file_columns: {e}")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.post("/api/extract")
async def smart_extraction(req: ExtractionRequest):
    """
    Smart Extraction: Extract specific sensors from multiple Excel/CSV files to new CSV/Excel.
    Consolidated implementation with robust filtering and structured output support.
    """
    try:
        session_id = req.session_id
        filenames = req.filenames
        required_sensors = req.sensors
        strict_mode = req.strict
        structured_mode = req.structured
        sheet_name = req.sheet_name
        selected_columns = req.columns
        section = req.section
        
        # If columns are explicitly selected, use them instead of sensors
        if selected_columns:
            required_sensors = selected_columns
            strict_mode = True # Use exact matches for explicit column selection
        
        logger.info(f"Add Metadata Request: {filenames}, sensors={required_sensors}, sheet={sheet_name}, section={section}")
        
        if not filenames:
             return JSONResponse({"success": False, "error": "No files provided"}, status_code=400)
        
        session_dir = storage.get_session_dir(session_id)
        
        # Validate files exist
        valid_paths = []
        for fname in filenames:
            fpath = session_dir / fname
            if fpath.exists():
                valid_paths.append(str(fpath))
            else:
                logger.warning(f"File not found: {fname}")
        
        if not valid_paths:
            return JSONResponse({"success": False, "error": "No valid files found"}, status_code=404)
            
        # PPTX Extraction Logic
        ppt_files = [f for f in valid_paths if f.lower().endswith('.pptx')]
        extracted_slides_info = []
        ppt_summary = ""
        ppt_note = None

        # Robust sensor parsing (handle strings with newlines/commas just in case)
        if isinstance(required_sensors, str):
            required_sensors = [s.strip() for s in re.split(r'[,\n\r]+', required_sensors) if s.strip()]
        elif isinstance(required_sensors, list):
            new_sensors = []
            for s in required_sensors:
                if isinstance(s, str) and ('\n' in s or ',' in s):
                    new_sensors.extend([x.strip() for x in re.split(r'[,\n\r]+', s) if x.strip()])
                else:
                    new_sensors.append(s)
            required_sensors = new_sensors

        # Filter out obvious log noise
        log_patterns = [
            r'INFO:', r'DEBUG:', r'WARNING:', r'ERROR:', r'CRITICAL:',
            r'HTTP/\d\.\d', r'GET /', r'POST /', r'PUT /', r'DELETE /',
            r'\d{4}-\d{2}-\d{2} \d{2}:\d{2}:\d{2}',
            r'\[\d{2}/[A-Z][a-z]{2}/\d{4}',
            r'127\.0\.0\.1',
            r' - - ',
        ]
        
        filtered_sensors = []
        noise_detected = False
        for s in required_sensors:
            is_noise = False
            for pattern in log_patterns:
                if re.search(pattern, s):
                    is_noise = True
                    noise_detected = True
                    break
            if not is_noise:
                filtered_sensors.append(s)
        
        if noise_detected:
            logger.info(f"Filtered log noise from sensors. Original: {len(required_sensors)}, Filtered: {len(filtered_sensors)}")
            
        required_sensors = filtered_sensors

        # Only return 400 if NO sensors AND NO PPT files
        if not required_sensors and not ppt_files:
            error_msg = "No valid sensors provided."
            if noise_detected:
                error_msg += " Input contained log data instead of sensor names."
            return JSONResponse({"success": False, "error": error_msg}, status_code=400)
        
        if ppt_files:
            logger.info(f"Processing PPTX extraction for {len(ppt_files)} files")
            
            # IMPROVED: Parse slide numbers from required_sensors (which acts as slide range for PPT)
            slide_numbers = []
            for s in required_sensors:
                try:
                    s_clean = str(s).strip()
                    if '-' in s_clean:
                        parts = s_clean.split('-')
                        if len(parts) == 2 and parts[0].isdigit() and parts[1].isdigit():
                            start, end = map(int, parts)
                            slide_numbers.extend(range(start, end + 1))
                    elif s_clean.isdigit():
                        slide_numbers.append(int(s_clean))
                except (ValueError, TypeError):
                    continue
            
            # IMPROVED: Default to all slides (or first 20) if no valid numbers found
            if not slide_numbers:
                logger.info("No slide numbers specified, defaulting to first 20 slides")
                slide_numbers = list(range(1, 21))
            
            for ppt_path in ppt_files:
                fname = Path(ppt_path).name
                slides = ppt_extractor.extract_slides(ppt_path, slide_numbers, output_dir=session_dir)
                
                for slide in slides:
                    if slide.get('image_path'):
                        # Add to report
                        storage.add_to_report_section(
                            session_id=session_id,
                            section=section or "Appendices",
                            item_type="image",
                            content=slide['image_path']
                        )
                        extracted_slides_info.append({
                            "filename": fname,
                            "slide": slide['slide_number'],
                            "title": slide['title']
                        })
            
            if extracted_slides_info:
                # Check extraction mode to inform user
                extraction_mode = extracted_slides_info[0].get('extraction_mode', 'full') if extracted_slides_info else 'unknown'
                
                if extraction_mode == 'text_only':
                    ppt_summary = f"Extracted {len(extracted_slides_info)} slides (TEXT ONLY) from {len(ppt_files)} PPTX files."
                    ppt_note = "To enable full PPT image export: pip install pywin32 && python venv/Scripts/pywin32_postinstall.py -install"
                else:
                    ppt_summary = f"Extracted {len(extracted_slides_info)} slides from {len(ppt_files)} PPTX files."
            else:
                logger.warning("No slides extracted from PPTX files")
                ppt_summary = "PPTX extraction attempted but no slides were found or extracted."

        # Run Sensor Extraction (for CSV/Excel)
        # Filter out PPTX files as they are handled separately
        data_files = [f for f in valid_paths if not f.lower().endswith(('.pptx', '.ppt'))]
        
        results, source_map = sensor_harvester.harvest_sensors(
            data_files, 
            required_sensors, 
            strict_mode=strict_mode, 
            structured_output=structured_mode, 
            orchestrator=enhanced_orchestrator,
            sheet_name=sheet_name
        )
        
        # If section is provided or we have a default, add to report
        target_section = section or "Analysis and Results"
        
        if target_section:
            if isinstance(results, pd.DataFrame) and not results.empty:
                summary_text = f"Extracted data from {', '.join(filenames)}"
                if sheet_name: summary_text += f" (Sheet: {sheet_name})"
                storage.add_to_report_section(session_id, target_section, "text", summary_text)
                # FIX: Pass structured data instead of markdown
                table_content = {
                    "columns": results.columns.tolist(),
                    "data": results.head(50).to_dict(orient='records')
                }
                storage.add_to_report_section(session_id, target_section, "table", table_content)
            elif isinstance(results, dict):
                for sheet_id, df in results.items():
                    if not df.empty:
                        storage.add_to_report_section(session_id, target_section, "text", f"Extracted from {sheet_id}")
                        # FIX: Pass structured data instead of markdown
                        table_content = {
                            "columns": df.columns.tolist(),
                            "data": df.head(50).to_dict(orient='records')
                        }
                        storage.add_to_report_section(session_id, target_section, "table", table_content)
        
        is_empty = False
        if isinstance(results, pd.DataFrame):
            is_empty = results.empty
        elif isinstance(results, dict):
            is_empty = not results
            
            # We don't return early here anymore to allow download_url generation
            # but we store the PPT summary to include in the final response
            pass
            
        if is_empty and not extracted_slides_info:
            missing = [s for s, src in source_map.items() if src == "NOT FOUND"]
            error_msg = "No data could be extracted."
            if ppt_files:
                error_msg += " PPT extraction failed (check if PowerPoint is installed)."
            if missing:
                error_msg += f" Sensors not found: {', '.join(missing[:5])}"
            return JSONResponse({"success": False, "error": error_msg}, status_code=400)

        # Save to file
        timestamp = datetime.now().strftime('%H%M%S')
        
        if structured_mode and isinstance(results, dict) and results:
            output_filename = f"Extracted_Structured_{timestamp}.xlsx"
            output_path = session_dir / output_filename
            
            with pd.ExcelWriter(output_path) as writer:
                for sheet_id, df in results.items():
                    sheet_name = sheet_id.split("::")[-1]
                    sheet_name = re.sub(r'[\\/*?\[\]]', '_', sheet_name)[:31]
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            first_sheet_name = list(results.keys())[0]
            master_df = results[first_sheet_name]
            row_count = sum(len(df) for df in results.values())
            summary = f"Extracted {row_count} rows across {len(results)} sheets from {len(valid_paths)} files."
        elif not isinstance(results, dict) or not results:
            # Handle case where results is empty but we have slides
            master_df = pd.DataFrame()
            output_filename = f"Extracted_Empty_{timestamp}.csv"
            output_path = session_dir / output_filename
            master_df.to_csv(output_path, index=False)
            summary = "No Excel data extracted."
        else:
            master_df = results
            output_filename = f"Extracted_Flat_{timestamp}.csv"
            output_path = session_dir / output_filename
            master_df.to_csv(output_path, index=False)
            summary = f"Extracted {len(master_df)} rows from {len(valid_paths)} files."

        storage.add_upload(session_id, output_filename)
        download_url = f"/workspace/{session_id}/{output_filename}"
        
        # Replace NaNs for JSON
        preview_data = master_df.head().replace({np.nan: None}).to_dict(orient='records')
            
        # Combine summaries
        final_summary = summary
        if ppt_summary:
            final_summary = f"{ppt_summary} {summary}".strip()

        return JSONResponse({
            "success": True,
            "download_url": download_url,
            "extracted_file": output_filename,
            "summary": final_summary,
            "preview": preview_data,
            "sensors_found": [s for s, src in source_map.items() if src != "NOT FOUND"],
            "extracted_slides": extracted_slides_info,
            "note": ppt_note
        })
            
    except Exception as e:
        logger.exception("Extraction error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)


@app.post("/api/columns/select")
async def save_column_selection(req: ColumnSelectionRequest):
    """Save column selection"""
    try:
        storage.set_column_selection(req.session_id, req.filename, req.x_column, req.y_columns)
        return JSONResponse({"success": True})
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.get("/api/columns/get")
async def get_column_selection(session_id: str = "default"):
    """Get saved column selection"""
    selection = storage.get_column_selection(session_id)
    return JSONResponse({"selection": selection})


@app.post("/api/chart/generate")
async def generate_chart_endpoint(req: ChartRequest):
    """
    Production-grade chart generation endpoint
    - Validates all inputs
    - Smart data sampling
    - Professional chart styling
    - Comprehensive error handling
    """
    try:
        session_dir = storage.get_session_dir(req.session_id)
        file_path = session_dir / req.filename
        
        # === STEP 1: File Validation ===
        if not file_path.exists():
            # Attempt recovery
            logger.warning(f"⚠️ File not found in session: {req.filename}. Attempting recovery...")
            found_path = storage.find_file(req.filename)
            
            if found_path and storage.copy_file_to_session(found_path, req.session_id):
                logger.info(f"✅ File recovered from: {found_path}")
            else:
                logger.error(f"❌ File not found: {file_path}")
                return JSONResponse({
                    "success": False,
                    "error": f"File '{req.filename}' not found in session (and could not be recovered)"
                }, status_code=404)
        
        logger.info(f"📂 Loading file: {file_path.name} ({file_path.stat().st_size / 1024:.1f} KB)")
        
        # === STEP 2: Load Data ===
        try:
            if req.filename.endswith('.csv'):
                df = pd.read_csv(str(file_path))
            elif req.filename.endswith(('.xlsx', '.xls')):
                # Use robust reader for Excel
                if req.sheet_name:
                    sheet_name = req.sheet_name
                else:
                    sheets = list_sheetnames(str(file_path))
                    sheet_name = sheets[0] if sheets else None
                
                if sheet_name:
                    df = read_excel_table(str(file_path), sheet_name)
                else:
                    df = pd.DataFrame()
            else:
                return JSONResponse({
                    "success": False,
                    "error": f"Unsupported file type: {req.filename}"
                }, status_code=400)
            
            logger.info(f"✅ Loaded {len(df)} rows × {len(df.columns)} columns (No sorting applied)")
            
        except Exception as e:
            logger.error(f"❌ File load error: {e}")
            return JSONResponse({
                "success": False,
                "error": f"Failed to parse file: {str(e)}"
            }, status_code=400)
        
        # === STEP 3: Column Validation ===
        # We check if the columns you asked for actually exist in the file.
        
        # Robust matching: Case-insensitive and whitespace-insensitive
        df.columns = [str(c).strip() for c in df.columns] # Clean columns
        
        def find_column_robust(target, columns):
            if target in columns:
                return target
            target_norm = target.lower().strip()
            for col in columns:
                if col.lower().strip() == target_norm:
                    return col
            return None

        missing_cols = []
        
        # Resolve X Column
        real_x = find_column_robust(req.x_column, df.columns)
        if real_x:
            req.x_column = real_x
        else:
            missing_cols.append(req.x_column)
            
        # Resolve Y Columns
        real_ys = []
        for col in req.y_columns:
            real_y = find_column_robust(col, df.columns)
            if real_y:
                real_ys.append(real_y)
            else:
                missing_cols.append(col)
        req.y_columns = real_ys
        
        if missing_cols:
            logger.error(f"❌ Missing columns: {missing_cols}")
            return JSONResponse({
                "success": False,
                "error": f"Columns not found in file: {', '.join(missing_cols)}",
                "available_columns": df.columns.tolist()
            }, status_code=400)
        
        # === STEP 4: Data Quality Check ===
        # We remove any rows that have empty values in the columns we are plotting.
        # This ensures the chart doesn't break or look weird.
        df_subset = df[[req.x_column] + req.y_columns].copy()
        df_clean = df_subset.dropna()
        
        if len(df_clean) == 0:
            return JSONResponse({
                "success": False,
                "error": "No valid data after removing null values"
            }, status_code=400)
        
        dropped_pct = ((len(df_subset) - len(df_clean)) / len(df_subset)) * 100
        if dropped_pct > 0:
            logger.warning(f"⚠️ Dropped {dropped_pct:.1f}% rows with null values")
        
        # === STEP 4.5: Run Advanced Analytics ===
        logger.info("🔍 Running advanced analytics...")
        advanced_analysis = analytics_engine.comprehensive_analysis(
            df=df_clean,
            x_col=req.x_column,
            y_cols=req.y_columns
        )
        
        # === STEP 5: Generate Interactive Plotly Chart ===
        # Use plotly for interactive charts with zoom, pan, hover
        output_path = str(session_dir / "chart")
        
        loop = asyncio.get_running_loop()
        success, msg, html_path, png_path, chart_id, summary, stats, plotly_json = await loop.run_in_executor(
            None,
            functools.partial(
                plotly_generator.generate_chart,
                df=df_clean,
                params={
                    "x_column": req.x_column,
                    "y_columns": req.y_columns,
                    "chart_type": req.chart_type,
                    "title": f"Multiple Sensors vs {req.x_column}" if len(req.y_columns) > 3 else f"{', '.join(req.y_columns)} vs {req.x_column}",
                    "data_points": len(df_clean),
                    "total_points": len(df_clean),
                    "user_query": req.user_query
                },
                output_path=output_path
            )
        )
        
        if success:
            # Prepare URLs for both HTML (interactive) and PNG (static for reports)
            html_filename = Path(html_path).name if html_path else None
            png_filename = Path(png_path).name if png_path else None
            
            chart_url_html = f"/workspace/{req.session_id}/{html_filename}" if html_filename else ""
            chart_url_png = f"/workspace/{req.session_id}/{png_filename}" if png_filename else ""
            
            # Primary URL for frontend (interactive)
            # Use PNG for chat display (<img> tag), HTML for interactive link if needed
            chart_url = chart_url_png
            
            logger.info(f"✅ Interactive chart generated: HTML={html_filename}, PNG={png_filename}")
            
            # === STEP 6: Generate Deep Insights (DISABLED - Pending better context) ===
            # try:
            #     deep_insights = enhanced_orchestrator.analyze_chart_deeply(
            #         chart_context={
            #             'chart_type': req.chart_type,
            #             'x_column': req.x_column,
            #             'y_columns': req.y_columns,
            #             'df_summary': {'total_points': len(df_clean)}
            #         },
            #         data_analysis=advanced_analysis
            #     )
            #     # Combine basic summary with deep insights
            #     enhanced_summary = f"{summary}\n\n{deep_insights}"
            # except Exception as e:
            #     logger.warning(f"Deep insights generation failed: {e}")
            #     enhanced_summary = summary
            
            enhanced_summary = summary
            
            # === STEP 7: Store Comprehensive Context ===
            logger.info(f"✅ Chart generated: {html_filename}")
            
            # Store in context manager
            context_manager.store_chart_context(
                session_id=req.session_id,
                chart_id=chart_id,
                filename=req.filename,
                x_column=req.x_column,
                y_columns=req.y_columns,
                chart_type=req.chart_type,
                df_summary={
                    'total_points': len(df_clean),
                    'columns': df_clean.columns.tolist()
                },
                advanced_analysis=advanced_analysis,
                chart_path=chart_url
            )
            
            # Save to history (legacy storage)
            storage.add_chart_to_history(req.session_id, {
                "chart_id": chart_id,
                "file": req.filename,
                "x_column": req.x_column,
                "y_columns": req.y_columns,
                "chart_type": req.chart_type,
                "data_points": len(df_clean),
                "summary": enhanced_summary,
                "stats": stats,
                "chart_path": png_path,  # Store PNG path for report generation
                "timestamp": datetime.now().isoformat()
            })
            
            # Append inference question
            final_summary = enhanced_summary + "\n\n**Do you want to add your inference or any observation about this data?**"
            
            return JSONResponse({
                "success": True,
                "chart_url": chart_url,
                "chart_id": chart_id,
                "summary": final_summary,
                "stats": stats,
                "advanced_analysis": {
                    "anomaly_count": sum(a.get('count', 0) for a in advanced_analysis.get('anomalies', {}).values()),
                    "recommendations": advanced_analysis.get('recommendations', [])[:3]
                },
                "metadata": {
                    "data_points_plotted": len(df_clean),
                    "total_data_points": len(df_clean),
                    "dropped_null_pct": round(dropped_pct, 1) if dropped_pct > 0 else 0
                },
                "plotly_json": plotly_json
            })
        else:
            logger.error(f"❌ Chart generation failed: {msg}")
            return JSONResponse({
                "success": False,
                "error": f"Chart generation failed: {msg}"
            }, status_code=500)
    
    except Exception as e:
        logger.exception("❌ Unexpected chart generation error")
        return JSONResponse({
            "success": False,
            "error": f"Unexpected error: {str(e)}"
        }, status_code=500)

@app.post("/api/heatmap/generate")
async def generate_heatmap_endpoint(req: HeatmapRequest):
    """Generate correlation heatmap"""
    try:
        session_dir = storage.get_session_dir(req.session_id)
        file_path = session_dir / req.filename
        
        if not file_path.exists():
            # Attempt recovery
            logger.warning(f"⚠️ File not found in session: {req.filename}. Attempting recovery...")
            found_path = storage.find_file(req.filename)
            
            if found_path and storage.copy_file_to_session(found_path, req.session_id):
                logger.info(f"✅ File recovered from: {found_path}")
            else:
                return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
            
        # Load data
        if req.filename.endswith('.csv'):
            df = pd.read_csv(str(file_path))
        elif req.filename.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(str(file_path))
        else:
            return JSONResponse({"success": False, "error": "Unsupported file type"}, status_code=400)
        
        # Filter columns if requested
        if req.columns:
            # Only keep columns that exist in the dataframe
            valid_cols = [c for c in req.columns if c in df.columns]
            if valid_cols:
                df = df[valid_cols]
        
        # Generate heatmap using Plotly generator
        output_path = str(session_dir / "heatmap")
        
        loop = asyncio.get_running_loop()
        success, msg, html_path, png_path, chart_id, summary, stats, plotly_json = await loop.run_in_executor(
            None,
            functools.partial(
                plotly_generator.generate_heatmap,
                df=df,
                output_path=output_path
            )
        )
        
        if success:
            # Prepare URLs
            html_filename = Path(html_path).name if html_path else None
            png_filename = Path(png_path).name if png_path else None
            
            heatmap_url = f"/workspace/{req.session_id}/{html_filename}" if html_filename else ""
            heatmap_image_url = f"/workspace/{req.session_id}/{png_filename}" if png_filename else ""
            
            # Save to history so feedback works
            storage.add_chart_to_history(req.session_id, {
                "chart_id": chart_id,
                "file": req.filename,
                "x_column": "Heatmap", # Placeholder
                "y_columns": req.columns if req.columns else df.columns.tolist(),
                "chart_type": "heatmap",
                "data_points": len(df),
                "summary": summary,
                "stats": stats,
                "chart_path": png_path,  # Store PNG path for report generation
                "timestamp": datetime.now().isoformat()
            })
            
            return JSONResponse({
                "success": True,
                "heatmap_url": heatmap_url,
                "heatmap_image_url": heatmap_image_url,
                "chart_id": chart_id,
                "analysis": summary,
                "plotly_json": plotly_json
            })
        else:
            return JSONResponse({"success": False, "error": msg}, status_code=500)
    
    except Exception as e:
        logger.exception("Heatmap error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

class CompareRequest(BaseModel):
    session_id: str
    files: List[str]
    columns: List[str]
    chart_type: str = "line"

@app.post("/api/compare/generate")
async def generate_comparison(req: CompareRequest):
    """Generate multi-file comparison chart using Plotly"""
    try:
        session_dir = storage.get_session_dir(req.session_id)
        datasets = {}
        
        # Load all files
        for fname in req.files:
            fpath = session_dir / fname
            if fpath.exists():
                try:
                    if fname.endswith('.csv'):
                        datasets[fname] = pd.read_csv(fpath)
                    else:
                        datasets[fname] = pd.read_excel(fpath)
                except Exception as e:
                    logger.warning(f"Failed to load {fname}: {e}")
        
        if not datasets:
            return JSONResponse({"success": False, "error": "No valid files found"}, status_code=400)
            
        # Auto-detect columns if not provided
        target_cols = req.columns
        if not target_cols:
            target_cols = comparative_analyzer.auto_detect_comparable_columns(datasets)
            if not target_cols:
                return JSONResponse({"success": False, "error": "Could not automatically detect comparable numeric columns. Please select columns manually."}, status_code=400)
        
        # Generate chart for the FIRST target column (for now, single chart)
        # TODO: Support multiple charts if multiple columns selected
        target_col = target_cols[0]
        
        chart_id = uuid.uuid4().hex[:8]
        output_path = session_dir / f"chart_{chart_id}.png" # Base path, will generate .html and .png
        
        success, html_path, png_path, plotly_json = comparative_analyzer.generate_comparison_chart(
            datasets=datasets,
            column=target_col,
            output_path=str(output_path),
            chart_type=req.chart_type
        )
        
        if success:
            # Generate stats and summary
            results = comparative_analyzer.compare_datasets(datasets, [target_col])
            
            # Use Agentic Insight
            summary = comparative_analyzer.generate_agentic_insight(enhanced_orchestrator, results)
            
            # URLs
            chart_url = f"/workspace/{req.session_id}/{Path(html_path).name}"
            chart_image_url = f"/workspace/{req.session_id}/{Path(png_path).name}" if png_path else None
            
            # Save to history
            storage.add_chart_to_history(req.session_id, {
                "chart_id": chart_id,
                "file": "Comparison",
                "x_column": "Time/Index",
                "y_columns": [target_col],
                "chart_type": req.chart_type,
                "data_points": sum(len(df) for df in datasets.values()),
                "summary": summary,
                "stats": results['statistics'],
                "chart_path": str(Path(png_path).absolute()) if png_path else None, # Store absolute PNG path
                "timestamp": datetime.now().isoformat()
            })
            
            return JSONResponse({
                "success": True,
                "chart_url": chart_url, # Interactive HTML
                "chart_image_url": chart_image_url, # Static PNG
                "chart_id": chart_id,
                "summary": summary,
                "stats": results['statistics'],
                "detected_columns": target_cols if not req.columns else [],
                "plotly_json": plotly_json
            })
        else:
            return JSONResponse({"success": False, "error": "Chart generation failed"}, status_code=500)
            
    except Exception as e:
        logger.exception("Comparison Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

class AdvancedCompareRequest(BaseModel):
    session_id: str
    base_file: str
    base_x_column: str
    comparison_files: List[Dict[str, str]] # List of {filename, y_column, label}

@app.post("/api/tools/compare")
async def generate_advanced_comparison(req: AdvancedCompareRequest):
    """
    Generate advanced multi-file comparison chart
    Matches the payload from comparison.js
    """
    try:
        session_dir = storage.get_session_dir(req.session_id)
        datasets = {}
        
        # 1. Load Base File
        base_path = session_dir / req.base_file
        if base_path.exists():
            if req.base_file.endswith('.csv'):
                datasets[req.base_file] = pd.read_csv(base_path)
            else:
                datasets[req.base_file] = pd.read_excel(base_path)
        else:
             return JSONResponse({"success": False, "error": f"Base file {req.base_file} not found"}, status_code=404)

        # 2. Load Comparison Files
        for item in req.comparison_files:
            fname = item['filename']
            if fname not in datasets:
                fpath = session_dir / fname
                if fpath.exists():
                    try:
                        if fname.endswith('.csv'):
                            datasets[fname] = pd.read_csv(fpath)
                        else:
                            datasets[fname] = pd.read_excel(fpath)
                    except Exception as e:
                        logger.warning(f"Failed to load {fname}: {e}")
        
        if not datasets:
            return JSONResponse({"success": False, "error": "No valid files found"}, status_code=400)
            
        # 3. Generate Chart
        chart_id = uuid.uuid4().hex[:8]
        output_path = session_dir / f"chart_{chart_id}.png"
        
        success, html_path, png_path, plotly_json = comparative_analyzer.generate_advanced_comparison_chart(
            datasets=datasets,
            series_config=req.comparison_files,
            base_x_column=req.base_x_column,
            output_path=str(output_path)
        )
        
        if success:
            # Generate summary using Agentic Insight if possible
            # For advanced comparison, we need to construct a results object manually or enhance the analyzer
            # For now, we'll construct a basic results object to pass to the orchestrator
            
            # Construct a synthetic results object for the orchestrator
            synthetic_results = {
                'dataset_names': list(datasets.keys()),
                'columns_compared': [s['y_column'] for s in req.comparison_files],
                'statistics': {}, # TODO: Calculate stats for advanced comparison
                'differences': {},
                'rankings': {}
            }
            
            # Calculate basic stats for the synthetic results
            for series in req.comparison_files:
                fname = series['filename']
                y_col = series['y_column']
                if fname in datasets and y_col in datasets[fname].columns:
                    data = pd.to_numeric(datasets[fname][y_col], errors='coerce').dropna()
                    if not data.empty:
                        synthetic_results['statistics'][f"{fname} - {y_col}"] = { # Use label as key
                             'mean': float(data.mean()),
                             'min': float(data.min()),
                             'max': float(data.max()),
                             'std': float(data.std())
                        }

            if enhanced_orchestrator:
                summary = enhanced_orchestrator.analyze_comparison(synthetic_results)
            else:
                summary = f"Comparison generated using base X-axis '{req.base_x_column}' from '{req.base_file}'.\n"
                summary += f"Compared {len(req.comparison_files)} series."
            
            # URLs
            chart_url = f"/workspace/{req.session_id}/{Path(html_path).name}"
            chart_image_url = f"/api/download/{Path(png_path).name}?session_id={req.session_id}" if png_path else None
            
            # Save to history
            storage.add_chart_to_history(req.session_id, {
                "chart_id": chart_id,
                "file": "Advanced Comparison",
                "x_column": req.base_x_column,
                "y_columns": [s['y_column'] for s in req.comparison_files],
                "chart_type": "line",
                "data_points": sum(len(df) for df in datasets.values()),
                "summary": summary,
                "stats": {},
                "timestamp": datetime.now().isoformat()
            })
            
            # Append inference question
            final_summary = summary + "\n\n**Do you want to add your inference or any observation about this data?**"
            
            return JSONResponse({
                "success": True,
                "chart_url": chart_url,
                "chart_image_url": chart_image_url,
                "chart_id": chart_id,
                "summary": final_summary,
                "plotly_json": plotly_json
            })
        else:
            return JSONResponse({"success": False, "error": "Comparison failed"}, status_code=500)
    except Exception as e:
        logger.exception("Advanced Comparison Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

class PPTExtractRequest(BaseModel):
    session_id: str
    filename: str
    slides: List[int]

@app.post("/api/ppt/extract")
async def extract_ppt_slides(req: PPTExtractRequest):
    """Extract slides and images from PPT"""
    try:
        session_dir = storage.get_session_dir(req.session_id)
        ppt_path = session_dir / req.filename
        
        if not ppt_path.exists():
            return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
            
        # We need to manually iterate to use our new image extractor
        # Or update ppt_extractor to handle it. 
        # Let's do it here for simplicity since we have the session_dir
        
        from pptx import Presentation
        prs = Presentation(str(ppt_path))
        results = []
        
        for slide_num in req.slides:
            idx = slide_num - 1
            if 0 <= idx < len(prs.slides):
                slide = prs.slides[idx]
                
                # Extract images
                prefix = f"slide_{slide_num}"
                images = ppt_extractor.extract_images_from_slide(slide, session_dir, prefix)
                
                # Add to report automatically? No, let user decide.
                # Just return the image paths (relative to workspace)
                
                rel_images = [f"/workspace/{req.session_id}/{Path(p).name}" for p in images]
                
                results.append({
                    "slide": slide_num,
                    "images": rel_images,
                    "text": ppt_extractor._get_slide_text(slide)
                })
                
        return JSONResponse({"success": True, "slides": results})
    except Exception as e:
        logger.exception("PPT Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

class ExcelSheetRequest(BaseModel):
    session_id: str
    filename: str

class ExcelExtractRequest(BaseModel):
    session_id: str
    filename: str
    sheet_name: str
    cell_range: Optional[str] = None  # e.g., "A1:D10"
    selected_columns: Optional[List[str]] = None

@app.post("/api/excel/sheets")
async def get_excel_sheets(req: ExcelSheetRequest):
    """Get list of sheets in an Excel file"""
    try:
        session_dir = storage.get_session_dir(req.session_id)
        file_path = session_dir / req.filename
        
        if not file_path.exists():
            return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
            
        if not req.filename.lower().endswith(('.xlsx', '.xls')):
             return JSONResponse({"success": False, "error": "Not an Excel file"}, status_code=400)

        xls = pd.ExcelFile(file_path)
        return JSONResponse({"success": True, "sheets": xls.sheet_names})
        
    except Exception as e:
        logger.exception("Excel Sheet List Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.post("/api/excel/extract")
async def extract_excel_table(req: ExcelExtractRequest):
    """Extract data from specific sheet with optional range/column filtering"""
    try:
        session_dir = storage.get_session_dir(req.session_id)
        file_path = session_dir / req.filename
        
        if not file_path.exists():
            return JSONResponse({"success": False, "error": "File not found"}, status_code=404)

        # Parse cell range if provided
        skiprows = None
        nrows = None
        usecols = None
        
        if req.cell_range:
            # Parse range like "A1:D10"
            try:
                import re
                match = re.match(r'([A-Z]+)(\d+):([A-Z]+)(\d+)', req.cell_range.upper())
                if match:
                    start_col, start_row, end_col, end_row = match.groups()
                    start_row = int(start_row)
                    end_row = int(end_row)
                    
                    # Convert column letters to indices
                    def col_to_num(col):
                        num = 0
                        for c in col:
                            num = num * 26 + (ord(c) - ord('A') + 1)
                        return num - 1
                    
                    start_col_idx = col_to_num(start_col)
                    end_col_idx = col_to_num(end_col)
                    
                    skiprows = start_row - 1 if start_row > 1 else 0
                    nrows = end_row - start_row + 1
                    usecols = list(range(start_col_idx, end_col_idx + 1))
                else:
                    return JSONResponse({"success": False, "error": "Invalid cell range format. Use format like 'A1:D10'"}, status_code=400)
            except Exception as e:
                logger.exception("Cell range parsing error")
                return JSONResponse({"success": False, "error": f"Cell range parsing error: {str(e)}"}, status_code=400)

        # Read specific sheet with optional range
        read_params = {"sheet_name": req.sheet_name}
        if skiprows is not None:
            read_params["skiprows"] = skiprows
        if nrows is not None:
            read_params["nrows"] = nrows
        if usecols is not None:
            read_params["usecols"] = usecols
            
        df = pd.read_excel(file_path, **read_params)
        
        # Filter columns if specified
        if req.selected_columns:
            available_cols = [col for col in req.selected_columns if col in df.columns]
            if available_cols:
                df = df[available_cols]
            else:
                return JSONResponse({"success": False, "error": "None of the selected columns found in sheet"}, status_code=400)
        
        # Replace NaN with empty string for better display
        df = df.fillna("")
        
        # Convert to records for table generation
        data = df.to_dict(orient='records')
        columns = df.columns.tolist()
        
        # Generate HTML preview
        html_preview = df.head(10).to_html(classes='table table-striped', index=False)
        
        return JSONResponse({
            "success": True, 
            "columns": columns,
            "data": data,
            "preview": html_preview,
            "row_count": len(df)
        })
        
    except Exception as e:
        logger.exception("Excel Extract Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.get("/api/chat/history")
async def get_chat_history(session_id: str = "default"):
    """Get chat history"""
    messages = storage.get_messages(session_id)
    return JSONResponse({"messages": messages})

@app.delete("/api/chat/history")
async def clear_chat_history(session_id: str = "default"):
    """Clear chat history"""
    storage.clear_messages(session_id)
    return JSONResponse({"success": True})

@app.get("/api/sessions")
async def get_sessions():
    """Get all chat sessions"""
    sessions = storage.get_all_sessions()
    return JSONResponse({"sessions": sessions})

@app.post("/api/chat")
async def intelligent_chat(req: ChatRequest):
    """
    Intelligent chat endpoint with deep context awareness and domain expertise
    """
    try:
        session_id = req.session_id
        user_message = req.message
        
        logger.info(f"💬 Chat query: {user_message[:100]}")
        
        # Get session context
        session_context = context_manager.get_context(session_id)
        
        # Get latest chart context if available
        latest_chart = context_manager.get_latest_chart_context(session_id)
        
        # Check if user is asking about specific data points
        query_lower = user_message.lower()
        data_query_response = None
        
        if latest_chart and any(word in query_lower for word in ['where', 'when', 'which points', 'what values']):
            # User might be asking a specific data query
            try:
                # Load the data file
                filename = latest_chart.get('filename')
                session_dir = storage.get_session_dir(session_id)
                file_path = session_dir / filename
                
                if file_path.exists():
                    if filename.endswith('.csv'):
                        df = pd.read_csv(file_path)
                    else:
                        df = pd.read_excel(file_path)
                    
                    # Try to answer the query
                    data_query_response = analytics_engine.query_data_context(
                        df=df,
                        x_col=latest_chart.get('x_column'),
                        y_cols=latest_chart.get('y_columns', []),
                        query=user_message
                    )
                    
                    if data_query_response and "couldn't parse" not in data_query_response.lower():
                        logger.info("✅ Answered data query directly")
                    else:
                        # Fallback to LLM if data query fails
                        data_query_response = None
            except Exception as e:
                logger.warning(f"Data query failed: {e}")
                data_query_response = None
        
        # Save user message
        storage.append_message(session_id, "user", user_message)
        
        # Get chat history for memory
        chat_history = storage.get_messages(session_id)
        
        # Try RAG
        rag_context = ""
        rag_sources = []
        confidence = 0.0
        rag_images = []
        
        try:
            from backend.tools.enhanced_rag_tools import rag_system
            if rag_system.initialized:
                rag_context, rag_sources, confidence, rag_images = rag_system.retrieve(user_message, top_k=3)
                logger.info(f"RAG: confidence={confidence:.2f}, images={len(rag_images)}")
        except Exception as e:
            logger.error(f"RAG retrieval error: {e}")

        # Use enhanced orchestrator for intelligent response
        if data_query_response:
            # We have a direct data answer
            response_text = data_query_response
        else:
            # Get deep, context-aware response
            advanced_analysis = None
            if latest_chart and 'advanced_analysis' in latest_chart:
                advanced_analysis = latest_chart['advanced_analysis']
            
            response_data = enhanced_orchestrator.get_enhanced_chat_response(
                user_query=user_message,
                session_context=session_context,
                chart_context=latest_chart,
                data_analysis=advanced_analysis,
                rag_context=rag_context,
                chat_history=chat_history
            )
            
            # Extract text from response (handle both dict and string)
            if isinstance(response_data, dict):
                response_text = response_data.get('content', str(response_data))
                
                # Handle actions
                if response_data.get('type') == 'agent_plan':
                    # Unified Agentic Flow: Use the agent platform
                    logger.info("🤖 Redirecting to Agent Platform for execution")
                    from backend.api.agent_endpoints import execute_agent_workflow, AgentExecuteRequest
                    
                    agent_req = AgentExecuteRequest(
                        message=user_message,
                        session_id=session_id,
                        auto_approve=True # Auto-approve for seamless chat experience
                    )
                    
                    agent_resp = await execute_agent_workflow(agent_req)
                    
                    # If it's a JSONResponse (error), extract content
                    if isinstance(agent_resp, JSONResponse):
                        import json
                        resp_content = json.loads(agent_resp.body.decode())
                        return JSONResponse({
                            "success": False,
                            "response": resp_content.get("response", "Agent execution failed")
                        })
                    
                    # Append the inference question if a chart was generated
                    final_response = agent_resp.response
                    
                    # Check for report artifact
                    report_url = None
                    download_ready = False
                    for artifact in agent_resp.artifacts:
                        if artifact.get('type') == 'report':
                            report_path = artifact.get('path')
                            if report_path:
                                report_url = f"/workspace/{session_id}/{Path(report_path).name}"
                                download_ready = True
                                break

                    if any(a.get('type') == 'chart' for a in agent_resp.artifacts):
                        final_response += "\n\n**Do you want to add your inference or any observation about this data?**"
                    
                    return JSONResponse({
                        "success": agent_resp.success,
                        "response": final_response,
                        "plan": agent_resp.plan,
                        "artifacts": agent_resp.artifacts,
                        "download_ready": download_ready,
                        "report_url": report_url
                    })

                elif response_data.get('type') == 'action':
                    action = response_data.get('action')
                    
                    if action == 'save_inference':
                        params = response_data.get('params', {})
                        storage.save_inference(
                            session_id=session_id,
                            content=params.get('content'),
                            section=params.get('section')
                        )
                        return JSONResponse({
                            "success": True,
                            "response": response_data.get('content')
                        })
                    
                    elif action == 'prompt_inference_section':
                        return JSONResponse({
                            "success": True,
                            "response": response_data.get('content')
                        })

                    elif action == 'generate_report':
                        logger.info("🚀 Triggering report generation from chat action")
                        
                        # Generate report
                        from backend.tools.report_generator import get_report_generator
                        report_gen = get_report_generator(llm_client)
                        output_path = storage.get_session_dir(session_id) / f"Thermal_Analysis_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                        
                        success, msg, report_path = report_gen.generate_report(
                            session_id=session_id,
                            storage=storage,
                            output_path=output_path
                        )
                        
                        if success:
                            # Standardize to /workspace/ URL which is known to work
                            report_url = f"/workspace/{session_id}/{Path(report_path).name}"
                            return JSONResponse({
                                "success": True,
                                "response": response_text + "\n\nReport generated successfully!",
                                "download_ready": True,
                                "report_url": report_url,
                                "filename": Path(report_path).name
                            })
                        else:
                            response_text += f"\n\n(Report generation failed: {msg})"
                    
                    elif action == 'plot_chart':
                        logger.info("📈 Triggering chart generation from chat action")
                        
                        # Try to get filename from context
                        filename = None
                        if latest_chart:
                            filename = latest_chart.get('filename')
                        elif session_context.get('files'):
                            filename = list(session_context.get('files').keys())[0]
                        
                        if filename:
                            try:
                                session_dir = storage.get_session_dir(session_id)
                                file_path = session_dir / filename
                                
                                if file_path.exists():
                                    if filename.endswith('.csv'):
                                        df = pd.read_csv(file_path)
                                    else:
                                        df = pd.read_excel(file_path)
                                    
                                    # Get numeric columns
                                    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
                                    
                                    if len(numeric_cols) >= 2:
                                        # INTELLIGENT COLUMN SELECTION
                                        # Check if user asked for "critical" parameters
                                        if any(word in user_message.lower() for word in ['critical', 'important', 'key', 'main']):
                                            logger.info("🔍 Using intelligent analyzer for critical parameter detection")
                                            # Use intelligent analyzer to find critical columns
                                            analysis = analytics_engine.analyze_dataframe(df)
                                            
                                            # Extract critical columns from analysis
                                            critical_cols = []
                                            if 'critical_columns' in analysis:
                                                critical_cols = analysis['critical_columns'][:3]  # Top 3
                                            elif 'basic_stats' in analysis:
                                                # Fallback: use columns with highest variance
                                                stats = analysis['basic_stats']
                                                sorted_cols = sorted(stats.items(), key=lambda x: x[1].get('std', 0), reverse=True)
                                                critical_cols = [col for col, _ in sorted_cols[:3]]
                                            
                                            if critical_cols:
                                                # Find time column for X-axis
                                                x_col = None
                                                for time_col in ['Time', 'Tm_Sampling', 'Timestamp', 'time', 'timestamp']:
                                                    if time_col in df.columns:
                                                        x_col = time_col
                                                        break
                                                
                                                if not x_col:
                                                    x_col = numeric_cols[0]  # Fallback to first numeric
                                                
                                                y_cols = [col for col in critical_cols if col != x_col][:3]
                                            else:
                                                # Fallback to first two numeric
                                                x_col = numeric_cols[0]
                                                y_cols = [numeric_cols[1]]
                                        else:
                                            # Default: use first two numeric columns
                                            x_col = numeric_cols[0]
                                            y_cols = [numeric_cols[1]]
                                        
                                        # Generate chart
                                        output_path = str(session_dir / "chart")
                                        
                                        loop = asyncio.get_running_loop()
                                        success, msg, html_path, png_path, chart_id, summary, stats, plotly_json = await loop.run_in_executor(
                                            None,
                                            functools.partial(
                                                plotly_generator.generate_chart,
                                                df=df,
                                                params={
                                                    "x_column": x_col,
                                                    "y_columns": y_cols,
                                                    "chart_type": response_data.get('params', {}).get('chart_type', 'line'),
                                                    "title": f"{', '.join(y_cols)} vs {x_col}",
                                                    "user_query": user_message
                                                },
                                                output_path=output_path
                                            )
                                        )
                                        
                                        if success:
                                            # Use PNG for chat display (<img> tag), HTML for interactive link if needed
                                            chart_url = f"/workspace/{session_id}/{Path(png_path).name}"
                                            
                                            # Append inference question
                                            final_response = response_text + f"\n\n✅ Generated chart: {', '.join(y_cols)} vs {x_col}"
                                            final_response += "\n\n**Do you want to add your inference or any observation about this data?**"
                                            
                                            return JSONResponse({
                                                "success": True,
                                                "response": final_response,
                                                "chart_url": chart_url,
                                                "chart_id": chart_id,
                                                "plotly_json": plotly_json
                                            })
                            except Exception as e:
                                logger.error(f"Auto-chart generation failed: {e}")
                                response_text += f"\n\n(Chart generation failed: {str(e)})"
            else:
                response_text = str(response_data)
        
        # Store query-response pair
        context_manager.add_query(session_id, user_message, response_text)
        
        # Save assistant response to storage (user message was saved at the start)
        storage.append_message(session_id, "assistant", response_text)
        
        logger.info(f"✅ Chat response generated ({len(response_text)} chars)")
        
        # Convert absolute paths to relative URLs for images
        formatted_images = []
        for img_path in (rag_images if 'rag_images' in locals() else []):
            try:
                # If it's already a URL or relative path, keep it
                if img_path.startswith('/') or img_path.startswith('http'):
                    formatted_images.append(img_path)
                    continue
                
                # Convert absolute path to relative URL
                path_obj = Path(img_path)
                if "assets" in path_obj.parts:
                    assets_idx = path_obj.parts.index("assets")
                    rel_path = "/".join(path_obj.parts[assets_idx:])
                    formatted_images.append(f"/{rel_path}")
                elif "workspace" in path_obj.parts:
                    ws_idx = path_obj.parts.index("workspace")
                    rel_path = "/".join(path_obj.parts[ws_idx:])
                    formatted_images.append(f"/{rel_path}")
                else:
                    formatted_images.append(img_path)
            except:
                formatted_images.append(img_path)
        
        return JSONResponse({
            "success": True,
            "response": response_text,
            "images": formatted_images
        })
        
    except Exception as e:
        logger.exception("Chat endpoint error")
        return JSONResponse({
            "success": False,
            "response": f"I encountered an error processing your query: {str(e)}",
            "error": str(e)
        }, status_code=500)


@app.post("/api/feedback")
async def submit_feedback(req: FeedbackRequest):
    """Handle chart feedback - mark as approved/rejected"""
    try:
        memory = storage.load_memory(req.session_id)
        chart_history = memory.get("chart_history", [])
        
        for chart in chart_history:
            if chart.get("chart_id") == req.chart_id:
                chart["feedback"] = req.feedback
                
                if req.feedback == "positive":
                    chart["approved"] = True
                    message = "✅ Chart approved! Select section and click '+ Report'"
                else:
                    chart["deleted"] = True
                    message = "🗑️ Chart removed"
                
                storage.save_memory(req.session_id, memory)
                return JSONResponse({"success": True, "message": message})
        
        return JSONResponse({"success": False, "error": "Chart not found"}, status_code=404)
    
    except Exception as e:
        logger.exception("Feedback error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)


class ReportAddRequest(BaseModel):
    session_id: str
    section: str
    item_type: str
    content: Any

@app.post("/api/report/add")
async def add_to_report(req: ReportAddRequest):
    """Add content to report section"""
    try:
        success, section = storage.add_to_report_section(
            req.session_id, 
            req.section, 
            req.item_type, 
            req.content
        )
        return JSONResponse({"success": success, "section": section})
    except Exception as e:
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

class ChartAddRequest(BaseModel):
    session_id: str
    chart_id: str
    chart_path: Optional[str] = None
    description: Optional[str] = None
    section: str = "Analysis and Results"  # Default section

@app.get("/api/download/{filename}")
async def download_file(filename: str, session_id: str):
    """Secure file download endpoint"""
    try:
        session_dir = storage.get_session_dir(session_id)
        file_path = session_dir / filename
        
        if not file_path.exists():
            return JSONResponse({"success": False, "error": "File not found"}, status_code=404)
        
        # Explicitly set Content-Disposition to force correct filename
        from fastapi.responses import Response
        with open(file_path, 'rb') as f:
            content = f.read()
        
        headers = {
            'Content-Disposition': f'attachment; filename="{filename}"'
        }
        
        return Response(
            content=content,
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document' if filename.endswith('.docx') else 'application/octet-stream',
            headers=headers
        )
    except Exception as e:
        logger.error(f"Download error: {e}")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.post("/api/report/add_chart")
async def add_chart_to_report(req: ChartAddRequest):
    """Add chart to specified report section"""
    try:
        # Get chart metadata from history
        memory = storage.load_memory(req.session_id)
        chart_history = memory.get("chart_history", [])
        
        chart_data = None
        for chart in chart_history:
            if chart.get("chart_id") == req.chart_id:
                chart_data = chart
                break
        
        if not chart_data:
            logger.error(f"Chart {req.chart_id} not found in history")
            return JSONResponse({"success": False, "error": "Chart not found in history"}, status_code=404)
        
        # Get the PNG path
        session_dir = storage.get_session_dir(req.session_id)
        chart_png_path = session_dir / f"chart_{req.chart_id}.png"
        
        # Use stored path if PNG doesn't exist at expected location
        if not chart_png_path.exists() and chart_data.get("chart_path"):
            chart_png_path = Path(chart_data["chart_path"])
        
        # Add chart to the selected section with complete metadata
        success, section = storage.add_to_report_section(
            req.session_id,
            req.section,
            "chart",
            {
                "chart_id": req.chart_id,
                "path": str(chart_png_path.absolute()),
                "x_column": chart_data.get("x_column", ""),
                "y_columns": chart_data.get("y_columns", []),
                "chart_type": chart_data.get("chart_type", "line"),
                "summary": req.description or chart_data.get("summary", "")
            }
        )
        
        if success:
            logger.info(f"✅ Chart {req.chart_id} added to {section} section")
        
        return JSONResponse({"success": success, "section": section})
    except Exception as e:
        logger.exception("Add chart to report error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

class AddTableRequest(BaseModel):
    session_id: str
    section: str
    table_data: Dict[str, Any]
    title: str

@app.post("/api/report/add_table")
async def add_table_to_report_endpoint(req: AddTableRequest):
    """Add a table to the report"""
    try:
        success, section = storage.add_to_report_section(
            req.session_id,
            req.section,
            "table",
            req.table_data
        )
        return JSONResponse({"success": success})
    except Exception as e:
        logger.exception("Add Table Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

@app.post("/api/report/generate")
async def generate_report_endpoint(session_id: str = Form("default")):
    """Generate comprehensive thermal analysis report"""
    try:
        from backend.tools.report_generator import get_report_generator
        
        # Get report generator
        report_gen = get_report_generator(llm_client)
        
        # Generate report
        output_path = storage.get_session_dir(session_id) / f"Thermal_Analysis_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        success, message, report_path = report_gen.generate_report(
            session_id=session_id,
            storage=storage,
            output_path=output_path
        )
        
        if success:
            report_url = f"/workspace/{session_id}/{Path(report_path).name}"
            
            return JSONResponse({
                "success": True,
                "message": message,
                "report_url": report_url,
                "filename": Path(report_path).name
            })
        else:
            return JSONResponse({
                "success": False,
                "error": message
            }, status_code=400)
    
    except Exception as e:
        logger.exception("Report generation error")
        return JSONResponse({
            "success": False,
            "error": str(e)
        }, status_code=500)






@app.post("/api/tools/compare")
async def advanced_comparison(req: Dict[str, Any]):
    """
    Advanced Comparison: Compare different columns across different files
    """
    try:
        session_id = req.get('session_id', 'default')
        base_file = req.get('base_file')
        base_x_column = req.get('base_x_column')
        comparison_files = req.get('comparison_files', []) # List of {filename, y_column, label}
        
        if not comparison_files:
            return JSONResponse({"success": False, "error": "No comparison series provided"}, status_code=400)
            
        session_dir = storage.get_session_dir(session_id)
        
        # Load all unique files involved
        unique_filenames = set()
        if base_file: unique_filenames.add(base_file)
        for item in comparison_files:
            unique_filenames.add(item['filename'])
            
        datasets = {}
        for fname in unique_filenames:
            fpath = session_dir / fname
            if fpath.exists():
                try:
                    if fname.lower().endswith('.csv'):
                        datasets[fname] = pd.read_csv(fpath)
                    else:
                        # Use robust reader for Excel
                        sheets = list_sheetnames(str(fpath))
                        sheet_name = sheets[0]
                        datasets[fname] = read_excel_table(str(fpath), sheet_name)
                except Exception as e:
                    logger.warning(f"Failed to load {fname} for comparison: {e}")
        
        if not datasets:
            return JSONResponse({"success": False, "error": "Could not load any datasets"}, status_code=400)
            
        # Generate advanced comparison chart
        chart_id = uuid.uuid4().hex[:8]
        chart_filename = f"AdvComparison_{chart_id}.png"
        output_path = session_dir / chart_filename
        
        success, html_path, png_path, plotly_json = comparative_analyzer.generate_advanced_comparison_chart(
            datasets=datasets,
            series_config=comparison_files,
            base_x_column=base_x_column,
            output_path=str(output_path)
        )
        
        if success:
            summary = f"Advanced comparison generated with {len(comparison_files)} series."
            
            # Save to history
            storage.add_chart_to_history(session_id, {
                "chart_id": chart_id,
                "file": "Advanced Comparison",
                "x_column": base_x_column or "Index",
                "y_columns": [s['y_column'] for s in comparison_files],
                "chart_type": "line",
                "summary": summary,
                "chart_path": png_path,
                "timestamp": datetime.now().isoformat()
            })

            return JSONResponse({
                "success": True,
                "chart_url": f"/workspace/{session_id}/{Path(html_path).name}",
                "chart_image_url": f"/workspace/{session_id}/{Path(png_path).name}" if png_path else None,
                "chart_id": chart_id,
                "summary": summary,
                "plotly_json": plotly_json
            })
        else:
            return JSONResponse({"success": False, "error": "Failed to generate advanced comparison chart"}, status_code=500)

    except Exception as e:
        logger.exception("Advanced Comparison Error")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)



@app.post("/api/compare")
async def compare_files(req: Dict[str, Any]):
    """
    File Comparison: Compare parameter across multiple files
    """
    try:
        session_id = req.get('session_id')
        filenames = req.get('filenames', [])
        parameter = req.get('parameter')
        
        if not filenames or len(filenames) < 2 or not parameter:
            return JSONResponse({"success": False, "error": "Need at least 2 files and 1 parameter"}, status_code=400)
            
        session_dir = storage.get_session_dir(session_id)
        
        # Load datasets
        datasets = {}
        for fname in filenames:
            fpath = session_dir / fname
            if fpath.exists():
                try:
                    if fname.endswith('.csv'):
                        datasets[fname] = pd.read_csv(fpath)
                    else:
                        datasets[fname] = pd.read_excel(fpath)
                except Exception as e:
                    logger.warning(f"Failed to load {fname}: {e}")
        
        if len(datasets) < 2:
             return JSONResponse({"success": False, "error": "Could not load enough files"}, status_code=400)
             
        # Generate comparison chart
        chart_id = uuid.uuid4().hex[:8]
        chart_filename = f"Comparison_{parameter}_{chart_id}.png"
        output_path = session_dir / chart_filename
        
        success, html_path, png_path, plotly_json = comparative_analyzer.generate_comparison_chart(
            datasets=datasets,
            column=parameter,
            output_path=str(output_path),
            chart_type='line'
        )
        
        if success:
            # Generate summary
            results = comparative_analyzer.compare_datasets(datasets, [parameter])
            summary = results.get('summary', f"Comparison generated for {parameter}")
            
            # Save to history
            storage.add_chart_to_history(session_id, {
                "chart_id": chart_id,
                "file": "Comparison",
                "x_column": "Time/Index",
                "y_columns": [parameter],
                "chart_type": "line",
                "data_points": sum(len(df) for df in datasets.values()),
                "summary": summary,
                "chart_path": png_path,
                "timestamp": datetime.now().isoformat()
            })

            return JSONResponse({
                "success": True,
                "chart_url": f"/workspace/{session_id}/{Path(html_path).name}",
                "chart_image_url": f"/workspace/{session_id}/{Path(png_path).name}" if png_path else None,
                "chart_id": chart_id,
                "summary": summary,
                "plotly_json": plotly_json
            })
        else:
            return JSONResponse({"success": False, "error": "Failed to generate comparison chart"}, status_code=500)

    except Exception as e:
        logger.error(f"Comparison failed: {e}")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

# Removed duplicate /api/tools/compare endpoint

@app.post("/api/inference/save")
async def save_inference(req: InferenceRequest):
    try:
        inference = storage.save_inference(
            session_id=req.session_id,
            content=req.content,
            chart_id=req.chart_id,
            section=req.section
        )
        return JSONResponse({
            "success": True,
            "inference": inference,
            "message": f"Inference saved successfully" + (f" to {req.section}" if req.section else "")
        })
    except Exception as e:
        logger.error(f"Error saving inference: {e}")
        return JSONResponse({"success": False, "error": str(e)}, status_code=500)

# ========== RUN SERVER ==========
if __name__ == "__main__":
    print("\n" + "="*50)
    print(" ThermoSense AI Server")
    print("="*50)
    print(f" URL: http://localhost:8000")
    print(f" Health: http://localhost:8000/health")
    print("="*50 + "\n")
    
    uvicorn.run(
        app,
        host="127.0.0.1",
        port=8000,
        log_level="info"
    )
